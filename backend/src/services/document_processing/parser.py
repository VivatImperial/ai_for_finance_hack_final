from __future__ import annotations

import asyncio
import io
from pathlib import Path
from typing import Any, Callable

import structlog
from docx import Document as DocxDocument
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation

from .models import MarkdownDocument


logger = structlog.get_logger(__name__)


class DocumentParser:
    """Lightweight parser that converts common office formats into Markdown."""

    async def parse(
        self, *, content_bytes: bytes, filename: str | None = None
    ) -> MarkdownDocument:
        return await asyncio.to_thread(
            self.parse_sync,
            content_bytes=content_bytes,
            filename=filename,
        )

    def parse_sync(
        self, *, content_bytes: bytes, filename: str | None = None
    ) -> MarkdownDocument:
        """Synchronous helper for scripts and tests."""
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop and loop.is_running():
            raise RuntimeError(
                "parse_sync cannot be used inside an active event loop; call `await parse(...)` instead."
            )

        markdown, metadata = self._parse_bytes(
            content_bytes=content_bytes, filename=filename
        )
        if not markdown.strip():
            raise RuntimeError("Parsed document is empty")

        return MarkdownDocument(content=markdown, metadata=metadata)

    # ------------------------------------------------------------------ #
    # Internal helpers
    # ------------------------------------------------------------------ #

    def _parse_bytes(
        self, *, content_bytes: bytes, filename: str | None
    ) -> tuple[str, dict[str, Any]]:
        extension = (Path(filename).suffix.lower() if filename else "") or ""
        parser = self._resolve_parser(extension)
        return parser(content_bytes, filename or "document")

    def _resolve_parser(
        self, extension: str
    ) -> Callable[[bytes, str], tuple[str, dict[str, Any]]]:
        if extension == ".pdf":
            return self._parse_pdf
        if extension in {".docx", ".dotx"}:
            return self._parse_docx
        if extension in {".pptx", ".ppsx"}:
            return self._parse_pptx
        return self._parse_plain_text

    def _parse_pdf(
        self, content_bytes: bytes, filename: str
    ) -> tuple[str, dict[str, Any]]:
        try:
            text = pdf_extract_text(io.BytesIO(content_bytes)) or ""
        except Exception as exc:
            raise RuntimeError("Failed to parse PDF document") from exc

        markdown = self._normalize_text(text)
        return markdown, self._base_metadata(filename)

    def _parse_docx(
        self, content_bytes: bytes, filename: str
    ) -> tuple[str, dict[str, Any]]:
        try:
            document = DocxDocument(io.BytesIO(content_bytes))
        except Exception as exc:
            raise RuntimeError("Failed to parse DOCX document") from exc

        lines: list[str] = []
        sections: list[dict[str, Any]] = []
        order = 0

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            style_name = paragraph.style.name if paragraph.style else ""
            heading_level = self._heading_level_from_style(style_name)

            if heading_level:
                level = min(heading_level, 6)
                lines.append(f"{'#' * level} {text}")
                sections.append({"title": text, "level": level, "order": order})
                order += 1
            else:
                lines.append(text)

        markdown = "\n\n".join(lines).strip()
        metadata = self._base_metadata(filename)
        metadata["sections"] = sections
        return markdown, metadata

    def _parse_pptx(
        self, content_bytes: bytes, filename: str
    ) -> tuple[str, dict[str, Any]]:
        try:
            presentation = Presentation(io.BytesIO(content_bytes))
        except Exception as exc:
            raise RuntimeError("Failed to parse PPTX document") from exc

        lines: list[str] = []
        sections: list[dict[str, Any]] = []

        for index, slide in enumerate(presentation.slides, start=1):
            title_shape = getattr(slide.shapes, "title", None)
            title_text = (
                title_shape.text.strip()
                if title_shape is not None and getattr(title_shape, "text", "")
                else f"Слайд {index}"
            )

            sections.append({"title": title_text, "level": 2, "order": index - 1})
            lines.append(f"## {title_text}")

            body_blocks: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if shape is title_shape:
                    continue

                paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs]
                text = "\n".join(filter(None, paragraphs)).strip()
                if text:
                    body_blocks.append(text)

            if body_blocks:
                lines.append("\n\n".join(body_blocks))

        markdown = "\n\n".join(lines).strip()
        metadata = self._base_metadata(filename)
        metadata["sections"] = sections
        return markdown, metadata

    def _parse_plain_text(
        self, content_bytes: bytes, filename: str
    ) -> tuple[str, dict[str, Any]]:
        text = self._decode_bytes(content_bytes)
        markdown = self._normalize_text(text)
        return markdown, self._base_metadata(filename)

    # ------------------------------------------------------------------ #
    # Utility helpers
    # ------------------------------------------------------------------ #

    def _base_metadata(self, filename: str) -> dict[str, Any]:
        return {"source_filename": filename, "sections": []}

    def _heading_level_from_style(self, style_name: str | None) -> int | None:
        if not style_name:
            return None

        normalized = style_name.lower()
        if "heading" not in normalized:
            return None

        digits = "".join(ch for ch in normalized if ch.isdigit())
        if digits.isdigit():
            return int(digits)
        return None

    def _decode_bytes(self, content_bytes: bytes) -> str:
        for encoding in ("utf-8", "utf-16", "cp1251", "latin-1"):
            try:
                return content_bytes.decode(encoding)
            except UnicodeDecodeError:
                continue

        return content_bytes.decode("utf-8", errors="ignore")

    def _normalize_text(self, text: str) -> str:
        lines = [line.rstrip() for line in text.splitlines()]
        normalized: list[str] = []
        blank = False

        for line in lines:
            stripped = line.strip()
            if not stripped:
                if not blank and normalized:
                    normalized.append("")
                    blank = True
                continue

            normalized.append(stripped)
            blank = False

        return "\n".join(normalized).strip()
